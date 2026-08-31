import os
import random
import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ==========================================
# 1. GENEREREN VAN DE VOORBEELDPLOTS
# ==========================================
os.makedirs("generated_plots", exist_ok=True)
sns.set_theme(style="white")

# Plot 1: Verdeling van Kaartwaarden (Simulatie 10.000 handen)
ranks = ['2', '3', '4', '5', '6', '7', '8', '9', '10', 'J', 'Q', 'K', 'A']
deck = [r for r in ranks for _ in range(4)]

random.seed(42)
drawn_ranks = []
pair_count = 0
total_hands = 10000

for _ in range(total_hands):
    hand = random.sample(deck, 2)
    drawn_ranks.extend(hand)
    if hand[0] == hand[1]:
        pair_count += 1

df_ranks = pd.DataFrame(drawn_ranks, columns=['Rank'])
rank_counts = df_ranks['Rank'].value_counts().reindex(ranks).reset_index()
rank_counts.columns = ['Rank', 'Count']

fig, ax = plt.subplots(figsize=(6, 3.8), dpi=200)
sns.barplot(
    data=rank_counts,
    x='Rank',
    y='Count',
    color='#1B4D3E',
    ax=ax,
    edgecolor='none',
)
ax.axhline(
    (total_hands * 2) / 13,
    color='#D9534F',
    linestyle='--',
    linewidth=1.5,
    label='Verwachte gemiddelde',
)
ax.set_title(
    'Distributie van 20.000 getrokken kaarten (10.000 handen)',
    fontsize=11,
    fontweight='bold',
    pad=10,
)
ax.set_xlabel('Kaartwaarde', fontsize=9)
ax.set_ylabel('Aantal keer getrokken', fontsize=9)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
ax.legend(frameon=False, fontsize=8)
plt.tight_layout()
plot1_path = "generated_plots/plot1_kaartverdeling.png"
plt.savefig(plot1_path)
plt.close()

# Plot 2: Pocket Pairs Frequentie (Empirisch vs Theoretisch)
empirical_pct = (pair_count / total_hands) * 100
theoretical_pct = (3 / 51) * 100  # ≈ 5.88%

fig, ax = plt.subplots(figsize=(5.5, 3.8), dpi=200)
categories = ['Theoretisch (3/51)', 'Jouw Simulatie (10k)']
values = [theoretical_pct, empirical_pct]
colors = ['#8A9BA8', '#2A6F97']

bars = ax.bar(categories, values, color=colors, width=0.5)
ax.set_ylabel('Percentage (%)', fontsize=9)
ax.set_ylim(0, 10)
ax.set_title(
    'Kans op een Pocket Pair (~5.88%)', fontsize=11, fontweight='bold', pad=10
)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)

for bar in bars:
    height = bar.get_height()
    ax.annotate(
        f'{height:.2f}%',
        xy=(bar.get_x() + bar.get_width() / 2, height),
        xytext=(0, 4),
        textcoords="offset points",
        ha='center',
        va='bottom',
        fontweight='bold',
        fontsize=9,
    )

plt.tight_layout()
plot2_path = "generated_plots/plot2_pair_frequentie.png"
plt.savefig(plot2_path)
plt.close()

# ==========================================
# 2. POWERPOINT GENERATOR (PYTHON-PPTX)
# ==========================================
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 Breedbeeld
prs.slide_height = Inches(7.5)

# Kleurenpalet
COLOR_BG = RGBColor(248, 249, 250)
COLOR_PRIMARY = RGBColor(18, 30, 49)
COLOR_ACCENT = RGBColor(27, 77, 62)
COLOR_MUTED = RGBColor(80, 90, 100)
COLOR_CARD = RGBColor(255, 255, 255)


def apply_background(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()
    return bg


def add_header(slide, title_text, category_text='WERKCOLLEGE 1 • DEEL 2'):
    header_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.1)
    )
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_ACCENT

    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY
    p_title.space_before = Pt(4)


# ------------------------------------------
# SLIDE 1: TITELSLIDE
# ------------------------------------------
slide_layout = prs.slide_layouts[6]
s1 = prs.slides.add_slide(slide_layout)

bg1 = s1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
)
bg1.fill.solid()
bg1.fill.fore_color.rgb = COLOR_PRIMARY
bg1.line.fill.background()

t_box = s1.shapes.add_textbox(
    Inches(1.2), Inches(2.2), Inches(11), Inches(3.2)
)
tf1 = t_box.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = 'POKERBOT ANALYTICS CHALLENGE'
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = RGBColor(79, 195, 161)

p = tf1.add_paragraph()
p.text = 'Week 1: De Eerste Deal 🃏'
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_before = Pt(8)

p = tf1.add_paragraph()
p.text = (
    'Van data genereren via kanssimulaties naar je allereerste beslis-bot in'
    ' Python'
)
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(200, 210, 220)
p.space_before = Pt(12)

s1.notes_slide.notes_text_frame.text = (
    "Welkom bij Deel 2! We gaan de zojuist geleerde basis van Python direct"
    " toepassen in een echte context: het simuleren van pokerkaarten en het"
    " bouwen van Bot Versie 1."
)

# ------------------------------------------
# SLIDE 2: WAAROM POKER?
# ------------------------------------------
s2 = prs.slides.add_slide(slide_layout)
apply_background(s2)
add_header(s2, 'Waarom Poker? De Ultieme Data Science Casus', 'WAAROM DEZE UITDAGING?')

col_width = Inches(3.7)
col_height = Inches(4.8)
col_y = Inches(1.8)

poker_reasons = [
    (
        Inches(0.8),
        '🎲 Onvolledige Informatie',
        'In schaken ligt alles open op tafel. Poker simuleert de echte praktijk:'
        ' beslissingen nemen onder onzekerheid met verborgen data.',
    ),
    (
        Inches(4.8),
        '📈 Expected Value (EV)',
        'Kansberekening regeert. Succes draait niet om één enkele winnende hand,'
        ' maar om beslissingen met een structureel positieve verwachtingswaarde.',
    ),
    (
        Inches(8.8),
        '🤖 Opponent Modeling',
        'Data science in actie: ontdek patronen in het gedrag van tegenstanders'
        ' en optimaliseer je logica om hun speelstijl te exploiteren.',
    ),
]

for x_pos, card_title, card_desc in poker_reasons:
    col = s2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, col_y, col_width, col_height
    )
    col.fill.solid()
    col.fill.fore_color.rgb = COLOR_CARD
    col.line.color.rgb = RGBColor(220, 225, 230)

    tf = col.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = Inches(0.3)

    p = tf.paragraphs[0]
    p.text = card_title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    p = tf.add_paragraph()
    p.text = card_desc
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_MUTED
    p.space_before = Pt(12)

s2.notes_slide.notes_text_frame.text = (
    "Benadruk dat poker een afspiegeling is van reële data science:"
    " besluitvorming onder onvolledige data en sturen op kansen op lange termijn."
)

# ------------------------------------------
# SLIDE 3: PLANNING & DEADLINE (13:40)
# ------------------------------------------
s3 = prs.slides.add_slide(slide_layout)
apply_background(s3)
add_header(s3, 'Programma & Planning voor vandaag')

card_left = s3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8),
    Inches(1.8),
    Inches(5.6),
    Inches(4.8),
)
card_left.fill.solid()
card_left.fill.fore_color.rgb = COLOR_CARD
card_left.line.color.rgb = RGBColor(220, 225, 230)

tf = card_left.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.4)
tf.margin_top = Inches(0.4)

p = tf.paragraphs[0]
p.text = '⏱️ Tijdverdeling'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

items_left = [
    (
        '15 min (Klassikaal nu)',
        'Concepten bespreken: simulaties, bot-functies & visual hierarchy.',
    ),
    (
        '35 min (Zelf aan de slag)',
        'Starten in het Jupyter Notebook en je bot-functie schrijven.',
    ),
    (
        'Huiswerk (tot wo 13:40)',
        'Afronden, visualisatie polijsten en inleveren via de API.',
    ),
]
for title, desc in items_left:
    p = tf.add_paragraph()
    p.text = f'• {title}'
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_ACCENT
    p.space_before = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = f'   {desc}'
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_MUTED

card_right = s3.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.8),
    Inches(1.8),
    Inches(5.7),
    Inches(4.8),
)
card_right.fill.solid()
card_right.fill.fore_color.rgb = COLOR_CARD
card_right.line.color.rgb = RGBColor(220, 225, 230)

tf = card_right.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.4)
tf.margin_top = Inches(0.4)

p = tf.paragraphs[0]
p.text = '🎯 De 3 Kerndoelen'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

items_right = [
    (
        '1. Data Genereren',
        '10.000 handen simuleren via random sampling zonder terugleggen.',
    ),
    (
        '2. Bot v1 Programmeren',
        'Een werkende Python-functie `kies_actie(beide_kaarten)` opleveren.',
    ),
    (
        '3. Visualiseren & Exporteren',
        'Verdeling plotten in Seaborn en als .py-bestand klaarzetten.',
    ),
]
for title, desc in items_right:
    p = tf.add_paragraph()
    p.text = f'• {title}'
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_ACCENT
    p.space_before = Pt(12)

    p2 = tf.add_paragraph()
    p2.text = f'   {desc}'
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_MUTED

s3.notes_slide.notes_text_frame.text = (
    "Let op de deadline: woensdag om 13:40 uur moet de API-check succesvol"
    " zijn doorlopen."
)

# ------------------------------------------
# SLIDE 4: DATA GENEREREN (MET PLOT 1)
# ------------------------------------------
s4 = prs.slides.add_slide(slide_layout)
apply_background(s4)
add_header(s4, 'Stap 1: 10.000 Handen Simuleren')

tb = s4.shapes.add_textbox(
    Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8)
)
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = 'Waarom data genereren via code?'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

bullets_s4 = [
    (
        'Monte Carlo Principe',
        'Door een willekeurig proces duizenden keren te herhalen, benaderen we'
        ' de werkelijke kansen.',
    ),
    (
        'Het Deck als Lijst',
        'Een lijst van 52 kaarten waaruit we telkens 2 kaarten trekken met'
        ' `random.sample()`.',
    ),
    (
        'Consistentie Check',
        'Elke kaartwaarde (2 t/m A) moet op de lange termijn even vaak'
        ' verschijnen (~7.69%).',
    ),
]
for b_title, b_desc in bullets_s4:
    p = tf.add_paragraph()
    p.text = f'• {b_title}: '
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_ACCENT
    p.space_before = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = f'  {b_desc}'
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_MUTED

s4.shapes.add_picture(
    plot1_path, Inches(6.6), Inches(2.0), width=Inches(5.8)
)

s4.notes_slide.notes_text_frame.text = (
    "Leg uit: we trekken kaarten ZONDER terugleggen binnen één hand. De rode"
    " stippellijn toont het theoretische gemiddelde."
)

# ------------------------------------------
# SLIDE 5: BOT V1 LOGICA (MET GEÜPDATETE CODE)
# ------------------------------------------
s5 = prs.slides.add_slide(slide_layout)
apply_background(s5)
add_header(s5, 'Stap 2: Wat is een Pokerbot in Python?')

card_top = s5.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8),
    Inches(1.8),
    Inches(11.7),
    Inches(1.6),
)
card_top.fill.solid()
card_top.fill.fore_color.rgb = COLOR_CARD
card_top.line.color.rgb = RGBColor(220, 225, 230)

tf = card_top.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = Inches(0.25)

p = tf.paragraphs[0]
p.text = 'Een bot is simpelweg een functie met een helder contract:'
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

p = tf.add_paragraph()
p.text = (
    '• INPUT: Twee gedeelde kaarten in een lijst `beide_kaarten` (bijv.'
    " `[('A', 'H'), ('K', 'D')]`)\n• LOGICA: Analyseer kaarten (in Week 1:"
    ' willekeurige geldige actie)\n• OUTPUT: Altijd exact één string:'
    " `'fold'`, `'call'` of `'raise'`"
)
p.font.size = Pt(11)
p.font.color.rgb = COLOR_MUTED
p.space_before = Pt(4)

code_box = s5.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.2)
)
code_box.fill.solid()
code_box.fill.fore_color.rgb = RGBColor(24, 30, 40)
code_box.line.fill.background()

tf = code_box.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = Inches(0.35)

code_text = (
    '# Voorbeeld: Bot Versie 1 (mijn_bot_week1.py)\n'
    'import random\n\n'
    'def kies_actie(beide_kaarten: list) -> str:\n'
    '    """\n'
    '    beide_kaarten: lijst van 2 kaarten, bijv. [(\'A\', \'H\'), (\'K\','
    " 'D')]\n"
    '    Geeft terug: \'call\', \'raise\', of \'fold\'\n'
    '    """\n'
    "    mogelijke_acties = ['fold', 'call', 'raise']\n"
    '    return random.choice(mogelijke_acties)  # Week 1: valide beslissing'
)
p = tf.paragraphs[0]
p.text = code_text
p.font.size = Pt(12)
p.font.name = 'Consolas'
p.font.color.rgb = RGBColor(120, 225, 175)

s5.notes_slide.notes_text_frame.text = (
    "Benadruk dat de functienaam 'kies_actie' en parameter 'beide_kaarten'"
    " exact zo moeten heten. De API test direct of deze functie bestaat en"
    " geldige acties returnt."
)

# ------------------------------------------
# SLIDE 6: VISUALISATIE & HIERARCHY (MET PLOT 2)
# ------------------------------------------
s6 = prs.slides.add_slide(slide_layout)
apply_background(s6)
add_header(
    s6, 'Stap 3: Visualiseren & Visual Hierarchy', 'DATA DESIGN & COMMUNICATIE'
)

tb = s6.shapes.add_textbox(
    Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8)
)
tf = tb.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = 'Maak grafieken die direct inzicht geven:'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

design_bullets = [
    (
        'Elimineer Ruis (Chart Junk)',
        'Geen overbodige 3D-effecten of zware zwarte rasterlijnen.',
    ),
    (
        'Doelgericht Kleurgebruik',
        'Gebruik kleur met een functie (bijv. accentkleur voor uitschieters of'
        ' benchmarks).',
    ),
    (
        'Duidelijke Actietitels',
        'Zet de conclusie in de titel: "Pocket Pairs vallen in ~5.88% van de'
        ' deals".',
    ),
]
for title, desc in design_bullets:
    p = tf.add_paragraph()
    p.text = f'🎨 {title}'
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_ACCENT
    p.space_before = Pt(10)

    p2 = tf.add_paragraph()
    p2.text = f'   {desc}'
    p2.font.size = Pt(11)
    p2.font.color.rgb = COLOR_MUTED

s6.shapes.add_picture(
    plot2_path, Inches(6.8), Inches(2.0), width=Inches(5.6)
)

s6.notes_slide.notes_text_frame.text = (
    "Verbind dit met Visual Hierarchy: laat zien hoe de datalabels direct"
    " vertellen wat belangrijk is."
)

# ------------------------------------------
# SLIDE 7: HET INLEVERPROCES & API (DEADLINE 13:40)
# ------------------------------------------
s7 = prs.slides.add_slide(slide_layout)
apply_background(s7)
add_header(s7, 'Het Eindproduct & Inleveren via de API')

col1 = s7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8),
    Inches(1.8),
    Inches(5.6),
    Inches(4.8),
)
col1.fill.solid()
col1.fill.fore_color.rgb = COLOR_CARD
col1.line.color.rgb = RGBColor(220, 225, 230)

tf = col1.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = Inches(0.4)

p = tf.paragraphs[0]
p.text = '📁 1. Het Python Bestand'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

p_body = (
    '• Sla je geteste functie op als:\n'
    '  `mijn_bot_week1.py`\n\n'
    '• Bevat géén onnodige print statements of oneindige loops.\n\n'
    '• Bevat de functie `kies_actie(beide_kaarten)`.'
)
p = tf.add_paragraph()
p.text = p_body
p.font.size = Pt(12)
p.font.color.rgb = COLOR_MUTED
p.space_before = Pt(10)

col2 = s7.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(6.8),
    Inches(1.8),
    Inches(5.7),
    Inches(4.8),
)
col2.fill.solid()
col2.fill.fore_color.rgb = COLOR_CARD
col2.line.color.rgb = RGBColor(220, 225, 230)

tf2 = col2.text_frame
tf2.word_wrap = True
tf2.margin_left = tf2.margin_top = Inches(0.4)

p = tf2.paragraphs[0]
p.text = '🚀 2. De Automatische API Check'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

p_body2 = (
    '• Voer de inlever-cel onderaan je notebook uit.\n\n'
    '• De API test direct:\n'
    '   ✔️ Draait de code zonder errors?\n'
    "   ✔️ Krijg je 'call', 'raise' of 'fold' terug?\n"
    '   ✔️ Bevat de grafiek een titel en labels?\n\n'
    '• ⏰ Deadline: Woensdag 13:40 uur.'
)
p = tf2.add_paragraph()
p.text = p_body2
p.font.size = Pt(12)
p.font.color.rgb = COLOR_MUTED
p.space_before = Pt(10)

s7.notes_slide.notes_text_frame.text = (
    "De deadline is woensdag om 13:40 uur. Zodra de API '200 OK - Inzending"
    " goedgekeurd' teruggeeft, is de afvink geregistreerd."
)

# ------------------------------------------
# SLIDE 8: AAN DE SLAG
# ------------------------------------------
s8 = prs.slides.add_slide(slide_layout)
apply_background(s8)
add_header(s8, 'Aan de slag! (35 minuten)', 'ZELFSTANDIG AAN HET WERK')

banner = s8.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8)
)
banner.fill.solid()
banner.fill.fore_color.rgb = COLOR_CARD
banner.line.color.rgb = RGBColor(200, 215, 225)

tf = banner.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_top = Inches(0.5)

p = tf.paragraphs[0]
p.text = '📋 Jouw Stappenplan voor nu in de klas:'
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

stappen = [
    (
        '1',
        'Open het Jupyter Notebook `Week1_Pokerbot_Opdracht.ipynb` in je'
        ' omgeving.',
    ),
    (
        '2',
        'Voer de simulatie-lus uit (10.000 trekkingen) en controleer of de data'
        ' klopt.',
    ),
    (
        '3',
        'Schrijf de `kies_actie(beide_kaarten)` functie voor Bot Versie 1.',
    ),
    (
        '4',
        'Maak de Seaborn verdelingsplot met een duidelijke titel en'
        ' accentkleur.',
    ),
    (
        '5',
        'Sla op als `mijn_bot_week1.py` en test de API-inlevercel onderaan.',
    ),
]

for nr, stap in stappen:
    p = tf.add_paragraph()
    p.text = f'{nr}.  {stap}'
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_PRIMARY
    p.space_before = Pt(10)

p_help = tf.add_paragraph()
p_help.text = (
    '⏰ Deadline: Woensdag 13:40 uur  |  🙋 Hulp nodig? Steek je hand op!'
)
p_help.font.size = Pt(12)
p_help.font.bold = True
p_help.font.color.rgb = COLOR_ACCENT
p_help.space_before = Pt(20)

s8.notes_slide.notes_text_frame.text = (
    "Zet eventueel een timer op het scherm van 35 minuten. Loop rond om"
    " studenten te helpen met vragen."
)

# Opslaan
pptx_filename = "Pokerbot_Upgrade_Week1.pptx"
prs.save(pptx_filename)
print(f"✅ Presentatie succesvol opgeslagen als: '{pptx_filename}' (8 slides)")
print(f"✅ Figuren opgeslagen in: '{plot1_path}' en '{plot2_path}'")